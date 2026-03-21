from ibm_watsonx_orchestrate.agent_builder.tools import tool
import io


@tool()
def my_tool(input: str) -> bytes:
    """Executes the tool's action based on the provided input.

    Args:
        input (str): The tool's input.

    Returns:
        bytes: The PPTX presentation as a byte sequence.
    """
    from pptx import Presentation

    #functionality of the tool
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Hello, World!"
    subtitle.text = f"Generated for {input}"

    # Save presentation to BytesIO and return as bytes
    pptx_bytes = io.BytesIO()
    presentation.save(pptx_bytes)
    pptx_bytes.seek(0)
    return pptx_bytes.getvalue()